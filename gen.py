import os

from reportlab.platypus import (
    SimpleDocTemplate,
    Paragraph,
    Spacer,
    PageBreak,
    Preformatted,
)
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.pagesizes import letter
from reportlab.graphics.shapes import Drawing, Rect, String, Line
from reportlab.lib import colors

from pptx import Presentation
from pptx.util import Inches, Pt


# ========= CONFIGURACIÓN DE SALIDA =========

OUTPUT_DIR = "./output"
os.makedirs(OUTPUT_DIR, exist_ok=True)

full_doc_path = os.path.join(OUTPUT_DIR, "sistema_seguridad_iot_full.pdf")
deploy_doc_path = os.path.join(OUTPUT_DIR, "sistema_seguridad_iot_deployment.pdf")
dev_doc_path = os.path.join(OUTPUT_DIR, "sistema_seguridad_iot_manual_desarrollador.pdf")
inst_doc_path = os.path.join(OUTPUT_DIR, "sistema_seguridad_iot_manual_instalador.pdf")
ppt_path = os.path.join(OUTPUT_DIR, "sistema_seguridad_iot_presentacion_ejecutiva.pptx")


# ========= ESTILOS REPORTLAB =========

styles = getSampleStyleSheet()
styles.add(ParagraphStyle(name="TitleCenter", parent=styles["Title"], alignment=1))
code_style = ParagraphStyle(
    "Code",
    parent=styles["Normal"],
    fontName="Courier",
    fontSize=8,
    leading=9,
)


# ========= 1) DOCUMENTO TÉCNICO INTEGRAL (FULL) =========

def build_full_document(path: str):
    story = []

    def add_title(txt):
        story.append(Paragraph(txt, styles["TitleCenter"]))
        story.append(Spacer(1, 20))

    def add_h2(txt):
        story.append(Paragraph(f"<b>{txt}</b>", styles["Heading2"]))
        story.append(Spacer(1, 6))

    def add_p(txt):
        # Soporta saltos de línea dobles como párrafos visualmente más separados
        txt = txt.replace("\n\n", "<br/><br/>")
        story.append(Paragraph(txt, styles["Normal"]))
        story.append(Spacer(1, 8))

    add_title("Documento Técnico Integral – Sistema de Seguridad IoT")

    # --- SRS ---
    add_h2("1. Introducción")
    add_p(
        "Este documento describe en detalle el diseño, la arquitectura, los requisitos y la operación "
        "de un sistema de Seguridad IoT basado en Raspberry Pi 4, ESP32 y cámaras, desplegado "
        "mediante contenedores Docker. Combina especificación de requisitos de software (SRS), "
        "documento de arquitectura (SAD) y guía operacional (DevOps)."
    )

    add_h2("2. Alcance")
    add_p(
        "El sistema debe integrar sensores PIR, cámaras, relés de iluminación y otros sensores IoT "
        "para proveer monitoreo y reacción automática ante eventos de movimiento en distintas zonas. "
        "La solución se ejecuta localmente pero está pensada para ser migrable a nubes como AWS IoT "
        "Core sin rediseñar la arquitectura central."
    )

    add_h2("3. Requisitos Funcionales")
    rf_text = """
RF1. El sistema debe recibir eventos de sensores PIR a través de MQTT.

RF2. El sistema debe asociar cada evento PIR a una zona lógica configurada en la base de datos.

RF3. El sistema debe activar automáticamente los relés de iluminación pertenecientes a la zona donde se produjo el evento PIR.

RF4. El sistema debe activar las cámaras asociadas a la misma zona para capturar imágenes o vídeo del evento.

RF5. Las cámaras deben poder ser activadas también de forma manual mediante API REST.

RF6. Las imágenes capturadas deben ser enviadas al backend y almacenadas como evidencias en disco y registradas en la base de datos.

RF7. Un servicio de IA debe procesar las evidencias y agregar metadatos (por ejemplo detección de personas).

RF8. El sistema debe registrar eventos y telemetría en PostgreSQL en formato JSONB.

RF9. El dashboard web debe permitir visualizar dispositivos, zonas, eventos recientes y evidencias, además de consultar mediciones históricas.

RF10. Todos los componentes lógicos (broker MQTT, backend REST, worker MQTT, IA, frontend y base de datos) deben poder ejecutarse como contenedores Docker.
"""
    add_p(rf_text)

    add_h2("4. Requisitos No Funcionales")
    rnf_text = """
RNF1. Disponibilidad: el sistema debe operar 24/7 con mínima intervención manual.

RNF2. Rendimiento: la latencia entre un evento PIR y la activación de relés no debe superar los 300 ms bajo carga nominal.

RNF3. Seguridad: el broker MQTT debe requerir autenticación, y la API REST debe estar protegida mediante autenticación basada en tokens.

RNF4. Escalabilidad: debe ser posible agregar nuevas zonas y dispositivos sin cambios en la arquitectura central.

RNF5. Portabilidad: toda la solución debe poder levantarse mediante Docker Compose en una Raspberry Pi 4, y en un futuro en una VM o servidor en la nube.

RNF6. Observabilidad: se deben generar logs de los servicios principales, con rotación controlada y posibilidad de inspección remota.

RNF7. Mantenibilidad: el código debe estar organizado por servicios (microservicios ligeros) con responsabilidades claras.
"""
    add_p(rnf_text)

    story.append(PageBreak())

    # --- Arquitectura lógica con diagrama de bloques ---
    add_h2("5. Arquitectura Lógica del Sistema")

    add_p(
        "La arquitectura se compone de varios bloques lógicos: dispositivos IoT (PIR, cámaras, relés), "
        "un broker MQTT, un worker MQTT que actúa como motor de reglas y ETL, un backend REST "
        "(Django/Flask), un servicio de IA para reconocimiento de imágenes, un frontend web y "
        "la base de datos PostgreSQL."
    )

    d = Drawing(500, 300)

    # Dispositivos IoT
    d.add(Rect(10, 200, 140, 80, strokeColor=colors.black, fillColor=colors.lightgrey))
    d.add(String(20, 240, "Dispositivos IoT", fontSize=10))
    d.add(String(20, 225, "PIR / Cámaras / Relés", fontSize=8))

    # Broker MQTT
    d.add(Rect(190, 200, 120, 80, strokeColor=colors.black, fillColor=colors.whitesmoke))
    d.add(String(200, 240, "Broker MQTT", fontSize=10))
    d.add(String(200, 225, "Mosquitto", fontSize=8))

    # Worker MQTT
    d.add(Rect(350, 200, 140, 80, strokeColor=colors.black, fillColor=colors.whitesmoke))
    d.add(String(360, 240, "Worker MQTT", fontSize=10))
    d.add(String(360, 225, "Reglas + ETL", fontSize=8))

    # Backend
    d.add(Rect(50, 80, 140, 80, strokeColor=colors.black, fillColor=colors.whitesmoke))
    d.add(String(60, 120, "Backend REST", fontSize=10))
    d.add(String(60, 105, "Django/Flask", fontSize=8))

    # Servicio IA
    d.add(Rect(220, 80, 120, 80, strokeColor=colors.black, fillColor=colors.whitesmoke))
    d.add(String(230, 120, "Servicio IA", fontSize=10))
    d.add(String(230, 105, "OpenCV + TFLite", fontSize=8))

    # BBDD
    d.add(Rect(380, 80, 100, 80, strokeColor=colors.black, fillColor=colors.whitesmoke))
    d.add(String(390, 120, "PostgreSQL", fontSize=10))

    # Flechas
    d.add(Line(150, 240, 190, 240))  # IoT -> MQTT
    d.add(Line(310, 240, 350, 240))  # MQTT -> Worker
    d.add(Line(420, 200, 420, 160))  # Worker -> DB
    d.add(Line(120, 160, 120, 80))   # Backend -> DB (conceptual)
    d.add(Line(280, 160, 280, 80))   # IA -> DB

    story.append(d)
    story.append(Spacer(1, 16))

    story.append(PageBreak())

    # --- Arquitectura Docker ---
    add_h2("6. Arquitectura de Contenedores Docker")

    add_p(
        "Cada componente lógico del sistema se despliega en un contenedor independiente. "
        "La Raspberry Pi 4 actúa como host Docker y orquesta los servicios mediante Docker Compose."
    )

    d2 = Drawing(500, 320)
    # Host
    d2.add(Rect(10, 40, 480, 260, strokeColor=colors.black, fillColor=None))
    d2.add(String(20, 290, "Host: Raspberry Pi 4 (Docker)", fontSize=10))

    containers = [
        ("mqtt", "Mosquitto", 30, 220),
        ("backend", "Django/Flask API", 180, 220),
        ("worker", "MQTT Worker", 330, 220),
        ("ai", "Servicio IA", 30, 130),
        ("frontend", "Dashboard Web", 180, 130),
        ("db", "PostgreSQL", 330, 130),
    ]

    for name, label, x, y in containers:
        d2.add(Rect(x, y, 120, 70, strokeColor=colors.darkblue, fillColor=colors.whitesmoke))
        d2.add(String(x + 10, y + 45, name, fontSize=9))
        d2.add(String(x + 10, y + 30, label, fontSize=8))

    story.append(d2)
    story.append(Spacer(1, 16))
    story.append(PageBreak())

    # --- ER simplificado ---
    add_h2("7. Diagrama Entidad–Relación (Simplificado)")

    add_p(
        "El siguiente diagrama en alto nivel muestra las relaciones principales entre usuarios, "
        "zonas, dispositivos, eventos, evidencias y mediciones."
    )

    d3 = Drawing(500, 320)

    # Users
    d3.add(Rect(30, 250, 100, 40, strokeColor=colors.black, fillColor=colors.whitesmoke))
    d3.add(String(40, 270, "users", fontSize=9))

    # Zones
    d3.add(Rect(200, 250, 100, 40, strokeColor=colors.black, fillColor=colors.whitesmoke))
    d3.add(String(210, 270, "zones", fontSize=9))

    # Devices
    d3.add(Rect(370, 250, 100, 40, strokeColor=colors.black, fillColor=colors.whitesmoke))
    d3.add(String(380, 270, "devices", fontSize=9))

    # Events
    d3.add(Rect(80, 150, 100, 40, strokeColor=colors.black, fillColor=colors.whitesmoke))
    d3.add(String(90, 170, "events", fontSize=9))

    # Evidences
    d3.add(Rect(220, 150, 100, 40, strokeColor=colors.black, fillColor=colors.whitesmoke))
    d3.add(String(230, 170, "evidences", fontSize=9))

    # Measurements
    d3.add(Rect(360, 150, 120, 40, strokeColor=colors.black, fillColor=colors.whitesmoke))
    d3.add(String(370, 170, "measurements", fontSize=9))

    # Relations
    d3.add(Line(130, 270, 200, 270))  # users - zones
    d3.add(Line(300, 270, 370, 270))  # zones - devices
    d3.add(Line(420, 250, 420, 190))  # devices - measurements
    d3.add(Line(130, 190, 80, 190))   # devices - events (aprox)
    d3.add(Line(300, 190, 220, 190))  # events - evidences

    story.append(d3)
    story.append(Spacer(1, 16))
    story.append(PageBreak())

    # --- Secuencia PIR ---
    add_h2("8. Diagrama de Secuencia: PIR → Luces y Cámaras")

    add_p(
        "Este diagrama representa el flujo de mensajes cuando un sensor PIR detecta movimiento. "
        "Se muestran las interacciones básicas entre el PIR, el broker MQTT, el worker, los relés "
        "de iluminación y las cámaras."
    )

    d4 = Drawing(500, 320)

    actors = ["PIR", "MQTT", "Worker", "Relé", "Cámara"]
    x_positions = [40, 130, 220, 310, 400]

    for x, label in zip(x_positions, actors):
        d4.add(String(x - 10, 300, label, fontSize=8))
        d4.add(Line(x, 290, x, 40))

    def msg(y, x1, x2, text):
        d4.add(Line(x1, y, x2, y))
        d4.add(String((x1 + x2) / 2 - 20, y + 5, text, fontSize=7))

    msg(260, x_positions[0], x_positions[1], "event/motion")
    msg(240, x_positions[1], x_positions[2], "mensaje MQTT")
    msg(210, x_positions[2], x_positions[3], "cmd/relay/set")
    msg(190, x_positions[2], x_positions[4], "cmd/capture_image")
    msg(160, x_positions[4], x_positions[2], "POST /frame")
    msg(130, x_positions[2], x_positions[1], "estado/ack")

    story.append(d4)
    story.append(Spacer(1, 16))

    add_h2("9. Operación General y Casos de Uso")
    add_p(
        "CASO DE USO 1: Detección automática.\n\n"
        "1. Un PIR detecta movimiento en la Zona A y publica un mensaje MQTT.\n"
        "2. El worker MQTT identifica la zona y registra un evento en la base de datos.\n"
        "3. El worker busca todos los dispositivos tipo 'relay' asociados a la zona y envía comandos para encender las luces.\n"
        "4. El worker también envía comandos a las cámaras de la zona para capturar imágenes o activar streaming.\n"
        "5. Las cámaras envían las imágenes al backend, donde se almacenan como evidencias.\n"
        "6. El servicio de IA analiza las evidencias y anota metadatos de detección.\n"
        "7. El dashboard actualiza la vista de eventos y muestra notificaciones al usuario."
    )

    add_p(
        "CASO DE USO 2: Visualización de telemetría.\n\n"
        "1. Un ESP32 de telemetría publica periódicamente datos en tópicos 'devices/&lt;id&gt;/telemetry'.\n"
        "2. El worker MQTT captura y almacena los datos en la tabla 'measurements'.\n"
        "3. El usuario abre el dashboard y consulta gráficos históricos de temperatura, humedad u otros parámetros.\n"
    )

    for i in range(10, 20):
        add_h2(f"10.{i} Detalles Adicionales de Implementación")
        add_p(
            "Esta sección describe lineamientos de calidad de código, manejo de errores, política de reconexión MQTT, "
            "y mejores prácticas para la implementación en Raspberry Pi. Los servicios deben manejar reconexiones "
            "automáticas al broker y a la base de datos, así como registrar adecuadamente los fallos de comunicación. "
            "Los contenedores se gestionan mediante Docker Compose, permitiendo reinicios automatizados y despliegues "
            "consistentes. Es recomendable utilizar archivos .env para parametrizar credenciales y endpoints."
        )

    doc_full = SimpleDocTemplate(path, pagesize=letter)
    doc_full.build(story)


# ========= 2) DOCUMENTO DE DEPLOYMENT =========

def build_deployment_document(path: str):
    story = []

    def add_title(txt):
        story.append(Paragraph(txt, styles["TitleCenter"]))
        story.append(Spacer(1, 20))

    def add_h2(txt):
        story.append(Paragraph(f"<b>{txt}</b>", styles["Heading2"]))
        story.append(Spacer(1, 6))

    def add_p(txt):
        story.append(Paragraph(txt.replace("\n\n", "<br/><br/>"), styles["Normal"]))
        story.append(Spacer(1, 8))

    add_title("Guía de Deployment – Sistema de Seguridad IoT")

    add_h2("1. Requisitos Previos")
    add_p(
        "• Raspberry Pi 4 con Raspberry Pi OS 64 bits.\n"
        "• Docker y Docker Compose instalados.\n"
        "• Usuario con permisos para ejecutar docker sin sudo.\n"
        "• Acceso a Internet para descargar imágenes base."
    )

    add_h2("2. Estructura de Docker Compose")
    add_p(
        "Un archivo docker-compose.yml orquesta todos los servicios: broker MQTT, backend REST, worker MQTT, "
        "servicio IA, frontend y base de datos PostgreSQL. Cada servicio tiene su propia imagen y configuración."
    )

    compose_example = """version: '3.8'
services:
  mqtt:
    image: eclipse-mosquitto:2
    container_name: mqtt
    ports:
      - "1883:1883"
    volumes:
      - ./mosquitto/conf:/mosquitto/config
      - ./mosquitto/data:/mosquitto/data

  db:
    image: postgres:14
    container_name: postgres
    environment:
      - POSTGRES_DB=seguridad
      - POSTGRES_USER=seguridad
      - POSTGRES_PASSWORD=seguridad123
    volumes:
      - ./postgres/data:/var/lib/postgresql/data

  backend:
    build: ./backend
    container_name: backend
    depends_on:
      - db
    environment:
      - DATABASE_URL=postgresql://seguridad:seguridad123@db:5432/seguridad
    ports:
      - "8000:8000"

  worker:
    build: ./worker
    container_name: worker
    depends_on:
      - mqtt
      - db

  ai:
    build: ./ai
    container_name: ai
    depends_on:
      - backend

  frontend:
    build: ./frontend
    container_name: frontend
    ports:
      - "3000:80"
"""
    story.append(Preformatted(compose_example, code_style))
    story.append(Spacer(1, 8))

    add_h2("3. CI/CD con GitHub Actions (Ejemplo)")
    ci_example = """name: CI

on:
  push:
    branches: [ main ]

jobs:
  build-and-test:
    runs-on: ubuntu-latest
    steps:
      - uses: actions/checkout@v4
      - name: Set up Docker Buildx
        uses: docker/setup-buildx-action@v3
      - name: Build backend image
        run: docker build -t backend:latest backend
      - name: Build worker image
        run: docker build -t worker:latest worker
      - name: Run unit tests
        run: |
          docker run --rm backend:latest pytest
"""
    story.append(Preformatted(ci_example, code_style))
    story.append(Spacer(1, 8))

    deploy_doc = SimpleDocTemplate(path, pagesize=letter)
    deploy_doc.build(story)


# ========= 3) MANUAL DEL DESARROLLADOR =========

def build_dev_manual(path: str):
    story = []

    def add_title(txt):
        story.append(Paragraph(txt, styles["TitleCenter"]))
        story.append(Spacer(1, 20))

    def add_h2(txt):
        story.append(Paragraph(f"<b>{txt}</b>", styles["Heading2"]))
        story.append(Spacer(1, 6))

    def add_p(txt):
        story.append(Paragraph(txt.replace("\n\n", "<br/><br/>"), styles["Normal"]))
        story.append(Spacer(1, 8))

    add_title("Manual del Desarrollador – Sistema de Seguridad IoT")

    add_h2("1. Estructura del Repositorio")
    add_p(
        "Se recomienda una estructura monorepo con subdirectorios para cada servicio:\n\n"
        "• backend/   – API REST (Django/Flask)\n"
        "• worker/    – Servicio MQTT → BD\n"
        "• ai/        – Servicio de IA con FastAPI\n"
        "• frontend/  – Dashboard web\n"
        "• infra/     – Archivos Docker Compose, configuración de Mosquitto, etc.\n"
        "• docs/      – Documentación técnica y diagramas\n"
    )

    add_h2("2. Convenciones de Código")
    add_p(
        "• Lenguaje principal: Python 3.10.\n"
        "• Estilo: PEP8, utilizando herramientas como black y isort.\n"
        "• Tests: pytest para cada servicio, con tests unitarios y de integración básicos.\n"
        "• Logs: usar logging estándar de Python, configurado vía dictConfig.\n"
    )

    add_h2("3. Backend REST")
    add_p(
        "El backend expone endpoints para gestionar dispositivos, zonas, eventos y evidencias. "
        "Se utiliza un ORM (Django ORM o SQLAlchemy) para interactuar con PostgreSQL. "
        "Se recomienda organizar el código en módulos: routers, models, schemas, services."
    )

    add_h2("4. Worker MQTT")
    add_p(
        "El worker se encarga de:\n"
        "• Suscribirse a topics MQTT de dispositivos ('devices/#' y 'security/#').\n"
        "• Parsear los tópicos para determinar tipo de mensaje (telemetría, evento, estado).\n"
        "• Insertar datos en las tablas 'measurements' y 'events'.\n"
        "• Implementar la lógica 'PIR → luces + cámaras' consultando la base de datos.\n"
    )

    add_h2("5. Servicio IA")
    add_p(
        "El servicio IA recibe rutas de evidencias (imágenes/vídeos) y devuelve metadatos de detección. "
        "Se puede implementar como un microservicio FastAPI con un endpoint '/analyze' que recibe un ID "
        "de evidencia, carga el archivo desde disco, ejecuta el modelo TFLite y actualiza la columna "
        "'ai_metadata' en la tabla 'evidences'."
    )

    add_h2("6. Flujo de Desarrollo")
    add_p(
        "1. Crear ramas por feature.\n"
        "2. Agregar tests unitarios.\n"
        "3. Ejecutar los tests localmente.\n"
        "4. Levantar el entorno Docker Compose en modo desarrollo.\n"
        "5. Probar interacciones end-to-end (PIR simulado → worker → backend → DB).\n"
    )

    dev_doc = SimpleDocTemplate(path, pagesize=letter)
    dev_doc.build(story)


# ========= 4) MANUAL DEL INSTALADOR / TÉCNICO =========

def build_installer_manual(path: str):
    story = []

    def add_title(txt):
        story.append(Paragraph(txt, styles["TitleCenter"]))
        story.append(Spacer(1, 20))

    def add_h2(txt):
        story.append(Paragraph(f"<b>{txt}</b>", styles["Heading2"]))
        story.append(Spacer(1, 6))

    def add_p(txt):
        story.append(Paragraph(txt.replace("\n\n", "<br/><br/>"), styles["Normal"]))
        story.append(Spacer(1, 8))

    add_title("Manual de Instalación y Puesta en Marcha – Sistema de Seguridad IoT")

    add_h2("1. Preparación de Hardware")
    add_p(
        "• Instalar la Raspberry Pi 4 con su fuente de poder adecuada.\n"
        "• Conectar red (WiFi o Ethernet) estable.\n"
        "• Conectar o preparar los dispositivos: ESP32 con PIR, ESP32-CAM o RPi Zero con cámara USB, "
        "módulos de relé con su alimentación aislada.\n"
    )

    add_h2("2. Instalación del Sistema Operativo")
    add_p(
        "• Grabar Raspberry Pi OS 64 bits en la tarjeta microSD.\n"
        "• Configurar el acceso SSH.\n"
        "• Actualizar paquetes ('sudo apt update && sudo apt upgrade').\n"
    )

    add_h2("3. Instalación de Docker y Docker Compose")
    add_p(
        "• Instalar Docker siguiendo la documentación oficial para ARM.\n"
        "• Instalar Docker Compose.\n"
        "• Agregar el usuario al grupo 'docker' y reiniciar sesión.\n"
    )

    add_h2("4. Despliegue de los Contenedores")
    add_p(
        "• Copiar el archivo docker-compose.yml al directorio de despliegue.\n"
        "• Ejecutar 'docker compose up -d'.\n"
        "• Verificar que los contenedores estén en ejecución con 'docker ps'.\n"
    )

    add_h2("5. Configuración de Dispositivos IoT")
    add_p(
        "• Configurar los ESP32 con las credenciales WiFi y la dirección IP de la Raspberry Pi.\n"
        "• Configurar el broker MQTT en los dispositivos (usuario, contraseña, puerto 1883).\n"
        "• Asignar un nombre de dispositivo coherente con la configuración de la base de datos.\n"
    )

    add_h2("6. Pruebas Iniciales")
    add_p(
        "• Verificar que los mensajes MQTT lleguen al broker usando herramientas como mqtt-explorer.\n"
        "• Confirmar que el worker registra eventos en la base de datos.\n"
        "• Activar manualmente un PIR y comprobar que se encienden las luces de la zona y se activan las cámaras.\n"
    )

    inst_doc = SimpleDocTemplate(path, pagesize=letter)
    inst_doc.build(story)


# ========= 5) PRESENTACIÓN EJECUTIVA (PPTX) =========

def build_presentation(path: str):
    prs = Presentation()

    # Slide 1 - título
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Sistema de Seguridad IoT"
    subtitle.text = "Arquitectura, Beneficios y Componentes"

    # Slide 2 - Arquitectura general
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "Arquitectura General"
    tf = body.text_frame
    tf.text = "Componentes principales:"
    p = tf.add_paragraph()
    p.text = "• Sensores PIR (ESP32)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Cámaras (ESP32-CAM / RPi Zero + USB)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Broker MQTT, Backend REST, Servicio IA, Dashboard, BD PostgreSQL"
    p.level = 1

    # Slide 3 - Arquitectura Docker
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "Arquitectura Docker"
    tf = body.text_frame
    tf.text = "Servicios en contenedores:"
    for txt in [
        "• mqtt – Mosquitto",
        "• backend – API REST (Django/Flask)",
        "• worker – Motor de reglas MQTT/BD",
        "• ai – Servicio de reconocimiento de imágenes",
        "• frontend – Dashboard Web",
        "• db – PostgreSQL",
    ]:
        p = tf.add_paragraph()
        p.text = txt
        p.level = 1

    # Slide 4 - Flujo PIR
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "Flujo Clave: PIR → Luces y Cámaras"
    tf = body.text_frame
    tf.text = "Resumen del flujo:"
    pasos = [
        "1. PIR detecta movimiento y envía mensaje MQTT.",
        "2. El worker procesa el evento y determina la zona.",
        "3. Se activan relés de iluminación y cámaras de la zona.",
        "4. Se capturan imágenes y se almacenan como evidencias.",
        "5. El servicio IA analiza las evidencias.",
    ]
    for s in pasos:
        p = tf.add_paragraph()
        p.text = s
        p.level = 1

    # Slide 5 - Beneficios
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "Beneficios Clave"
    tf = body.text_frame
    tf.text = "• Arquitectura modular y escalable."
    p = tf.add_paragraph()
    p.text = "• Operación totalmente local, con posibilidad de migrar a la nube."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Integración flexible de distintos tipos de dispositivos."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Historial completo de eventos y evidencias."
    p.level = 1

    prs.save(path)


# ========= MAIN =========

if __name__ == "__main__":
    print("Generando documentos en:", OUTPUT_DIR)
    build_full_document(full_doc_path)
    print(" - Documento técnico integral OK")
    build_deployment_document(deploy_doc_path)
    print(" - Guía de deployment OK")
    build_dev_manual(dev_doc_path)
    print(" - Manual del desarrollador OK")
    build_installer_manual(inst_doc_path)
    print(" - Manual del instalador OK")
    build_presentation(ppt_path)
    print(" - Presentación ejecutiva OK")
    print("Listo.")
